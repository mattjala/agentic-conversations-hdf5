"""Build agentic-conversations-hdf5-week3.pptx.

Covers the HDF5 schema optimizations since 2026-05-12, the benchmark
methodology (to pre-empt 'imperfect measure' concerns), the per-change impact
table with applicability footnotes for the live-session-hijack case, and the
live Claude Code session capture feature.

Figures come from benchmarks/opt_figures.py (run that first).
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
FIGS = ROOT / "figures" / "slides"
OUT = ROOT / "results" / "agentic-conversations-hdf5-week3.pptx"

# ── Palette (matches prior decks) ────────────────────────────────────────────
BG = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT = RGBColor(0x16, 0x21, 0x3e)
HIGHLIGHT = RGBColor(0x0f, 0x3c, 0x96)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GREY = RGBColor(0xcc, 0xd6, 0xe0)
YELLOW = RGBColor(0xff, 0xd1, 0x66)
DIM_WHITE = RGBColor(0xa0, 0xb4, 0xc8)
GREEN = RGBColor(0x6f, 0xc2, 0x76)
RED = RGBColor(0xe8, 0x6a, 0x6a)
CODEBG = RGBColor(0x0d, 0x1b, 0x2a)

W = Inches(13.333)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, text, left, top, width, height, size=18, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_para(tf, text, size=16, bold=False, color=WHITE, indent=0,
             align=PP_ALIGN.LEFT, space_before=3):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    if indent:
        p.level = indent
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def title_bar(slide, title, subtitle=""):
    add_rect(slide, 0, 0, W, Inches(1.05), HIGHLIGHT)
    add_text(slide, title, Inches(0.35), Inches(0.12), Inches(12.6), Inches(0.55),
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.35), Inches(0.66), Inches(12.6), Inches(0.34),
                 size=14, color=LIGHT_GREY)


def add_figure(slide, name, left, top, width, height):
    p = FIGS / f"{name}.png"
    if p.exists():
        slide.shapes.add_picture(str(p), left, top, width=width, height=height)
    else:
        add_text(slide, f"[missing figure: {name}]", left, top, width, Inches(0.4),
                 size=12, color=RED)


def bullets(slide, items, left, top, width, height, head=None, head_color=YELLOW,
            size=16):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = tf.paragraphs[0]
    if head:
        first.add_run().text = head
        first.runs[0].font.size = Pt(18)
        first.runs[0].font.bold = True
        first.runs[0].font.color.rgb = head_color
        start = None
    else:
        start = items[0]
        first.add_run().text = ("• " + start[0]) if isinstance(start, tuple) else ("• " + start)
        first.runs[0].font.size = Pt(size)
        first.runs[0].font.color.rgb = LIGHT_GREY
    for it in (items if head else items[1:]):
        txt, ind, col = (it if isinstance(it, tuple) else (it, 0, LIGHT_GREY))
        prefix = "  " * ind + ("– " if ind else "• ")
        add_para(tf, prefix + txt, size=size - (1 if ind else 0), color=col, indent=0)
    return txb


def add_table(slide, rows, left, top, width, height, col_widths=None, font=11):
    nrows, ncols = len(rows), len(rows[0])
    tbl = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
    tbl.first_row = False
    tbl.horz_banding = False
    for r in range(nrows):
        for c in range(ncols):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = HIGHLIGHT if r == 0 else (
                ACCENT if r % 2 else RGBColor(0x10, 0x18, 0x2e))
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(rows[r][c])
            run.font.size = Pt(font + 1 if r == 0 else font)
            run.font.bold = (r == 0)
            run.font.color.rgb = WHITE if r == 0 else LIGHT_GREY
    if col_widths:
        for c, wd in enumerate(col_widths):
            tbl.columns[c].width = wd
    return tbl


# ════════════════════════════════════════════════════════════════════════════
prs = new_prs()

# ── 1. Title ─────────────────────────────────────────────────────────────────
s = blank_slide(prs); fill_bg(s)
add_rect(s, 0, Inches(2.3), W, Inches(2.9), HIGHLIGHT)
add_text(s, "HDF5 Schema Optimizations + Live Session Capture",
         Inches(0.6), Inches(2.45), Inches(12.1), Inches(1.2),
         size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Progress since 2026-05-12  ·  methodology, per-change impact, and the live Claude Code hook",
         Inches(0.6), Inches(3.7), Inches(12.1), Inches(0.7),
         size=18, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
add_text(s, "Matthew Larson · The HDF Group",
         Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5),
         size=15, color=DIM_WHITE, align=PP_ALIGN.CENTER)

# ── 2. What changed since 5/12 ───────────────────────────────────────────────
s = blank_slide(prs); fill_bg(s)
title_bar(s, "What changed since 5/12", "Optimizations to the HDF5 conversation-log backend")
bullets(s, [
    ("Batched-write buffer — new `batch_size` knob collapses ~9 per-row HDF5 writes into one block write per batch", 0, LIGHT_GREY),
    ("Up to ~60× text write throughput at batch=100; batch=1 = old behavior", 1, DIM_WHITE),
    ("Embeddings consolidated into one (N, dim) dataset + has_embedding flag", 0, LIGHT_GREY),
    ("Full-scan read drops from 10,005 HDF5 calls to 8", 1, DIM_WHITE),
    ("Per-op profiler — exact HDF5 call counts per dataset, per phase", 0, LIGHT_GREY),
    ("ORC backend evaluated as an alternative columnar log format", 0, LIGHT_GREY),
    ("Live Claude Code session hook — incremental JSONL → HDF5 as you work", 0, GREEN),
], Inches(0.6), Inches(1.3), Inches(12.2), Inches(5.8), head="What changed")

# ── 3. Methodology I ─────────────────────────────────────────────────────────
s = blank_slide(prs); fill_bg(s)
title_bar(s, "How the benchmark works", "Methodology — what is and isn't being measured")
bullets(s, [
    ("Synthetic, in-memory turns. Text is drawn from a fixed 10,000-string pool; no source file is read during timing.", 0, LIGHT_GREY),
    ("→ measures backend write/read cost, not JSON parsing or disk-source I/O", 1, DIM_WHITE),
    ("Two scenarios:", 0, LIGHT_GREY),
    ("text-only — no embeddings/arrays; proxy for real Claude Code logs", 1, DIM_WHITE),
    ("array-heavy — 1536-d embedding + 1000-float tool result per turn; vector-store stress test", 1, DIM_WHITE),
    ("Metrics:", 0, LIGHT_GREY),
    ("write t/s = N / wall-time of the append loop incl. close()", 1, DIM_WHITE),
    ("ctx ms = cold-open + get_recent_context(20), median of 5 (a fresh consumer)", 1, DIM_WHITE),
    ("scan = full read of all N; size = bytes on disk", 1, DIM_WHITE),
    ("Sizes N = 50 / 200 / 1,000 / 5,000 turns (5,000 is the reported headline)", 0, LIGHT_GREY),
], Inches(0.6), Inches(1.3), Inches(12.2), Inches(5.8), head="The harness (benchmark.py)")

# ── 4. Methodology II — rigor & caveats ──────────────────────────────────────
s = blank_slide(prs); fill_bg(s)
title_bar(s, "Rigor and honest caveats", "Why the numbers are trustworthy — and where they aren't")
bullets(s, [
    ("Exact call counts. The per-op profiler wraps h5py to count every read/write/resize/create per dataset — these are counted, not estimated.", 0, GREEN),
    ("Caveat: the profiler's per-op milliseconds are inflated by its own wrapping. We report call counts, not those times.", 0, RED),
    ("Caveat: run-to-run variance ≈ 10% (machine load). We report relative speedups, not absolute t/s.", 0, RED),
    ("Caveat: the synthetic text pool is repetitive, which flatters gzip slightly — affects file size, not throughput comparisons.", 0, RED),
], Inches(0.6), Inches(1.3), Inches(12.2), Inches(5.8), head="Reading the results responsibly", size=15)

# ── 5. Impact table ──────────────────────────────────────────────────────────
s = blank_slide(prs); fill_bg(s)
title_bar(s, "Per-change impact", "Each optimization and its measured effect (text-only, N=5000 unless noted)")
rows = [
    ["Change", "Metric", "Before → After", "Live hook?"],
    ["libver='latest'", "write t/s", "296 → 319", "yes"],
    ["Pre-allocation (double)", "write t/s", "319 → 409", "yes"],
    ["Core VFD", "write t/s", "409 → 380 (rejected)", "no ¹"],
    ["zlib-ng + HDF5 2.0", "write t/s / ctx", "409 → 489 / −16%", "yes ²"],
    ["Chunk tuning", "—", "no change (defaults best)", "yes"],
    ["Batched buffer (b=100)", "write t/s", "489 → 25,870", "no ³"],
    ["Embedding consolidation", "array read calls", "10,005 → 8", "n/a ⁴"],
]
add_table(s, rows, Inches(0.5), Inches(1.25), Inches(12.3), Inches(4.4),
          col_widths=[Inches(3.4), Inches(3.0), Inches(3.7), Inches(2.2)], font=13)
add_text(s, "Live hook? = does the gain transfer to the shipped live-session capture path (see next slide for footnotes).",
         Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5), size=12, color=DIM_WHITE)

# ── 6. Footnotes: applicability to the live hijack ───────────────────────────
s = blank_slide(prs); fill_bg(s)
title_bar(s, "Does it apply to the live hook?", "Footnotes — where benchmark gains do / don't transfer to real Claude Code capture")
bullets(s, [
    ("¹ Core VFD was evaluated and rejected (regressions on array sessions); it is not in any shipped path.", 0, LIGHT_GREY),
    ("² zlib-ng compression gains depend on content. The synthetic pool is repetitive (flatters gzip); real Claude Code text still compresses well, but exact size deltas will differ.", 0, LIGHT_GREY),
    ("³ The batched buffer accelerates BULK conversion. The live hook syncs incrementally — only the few new JSONL lines per UserPromptSubmit/Stop event — so there is little to batch. Big win applies to `convert`, not the per-event hook.", 0, YELLOW),
    ("⁴ Embedding consolidation helps vector-store / array workloads. Claude Code logs carry NO embeddings, so the live hijack is effectively the text-only path; this gain is not exercised.", 0, YELLOW),
], Inches(0.6), Inches(1.3), Inches(12.2), Inches(5.8), head="Honest scope of each result", size=15)

# ── 7. Batched buffer figure ─────────────────────────────────────────────────
s = blank_slide(prs); fill_bg(s)
title_bar(s, "Batched-write buffer", "Collapsing ~9 per-row writes into one block write per batch")
add_figure(s, "opt_batch_scaling", Inches(0.6), Inches(1.3), Inches(7.4), Inches(4.3))
bullets(s, [
    ("Mechanism: HDF5 calls drop from ~9N to ~9·(N/batch)", 0, LIGHT_GREY),
    ("Text: 414 → 25,870 t/s (batch 1 → 100) ≈ 62×", 0, GREEN),
    ("Array gains less (≈2.4×): per-turn cost is embedding/tool writes, not row writes", 0, DIM_WHITE),
    ("File size and read latency unchanged — pure write-path win", 0, DIM_WHITE),
    ("batch=1 == old behavior (no overhead)", 0, DIM_WHITE),
], Inches(8.2), Inches(1.4), Inches(4.7), Inches(5.4), head="Takeaways", size=14)

# ── 8. Cumulative ────────────────────────────────────────────────────────────
s = blank_slide(prs); fill_bg(s)
title_bar(s, "Cumulative effect", "Stacking the write-path optimizations (text-only, N=5000)")
add_figure(s, "opt_cumulative", Inches(0.6), Inches(1.3), Inches(7.8), Inches(4.5))
bullets(s, [
    ("Each bar adds one optimization", 0, LIGHT_GREY),
    ("Note the log scale", 0, DIM_WHITE),
    ("libver/pre-alloc/zlib-ng: ~1.65× combined", 0, DIM_WHITE),
    ("The batched buffer is the step change", 0, GREEN),
    ("≈87× over the original baseline", 0, GREEN),
], Inches(8.6), Inches(1.4), Inches(4.3), Inches(5.4), head="Takeaways", size=14)

# ── 9. Embedding consolidation ───────────────────────────────────────────────
s = blank_slide(prs); fill_bg(s)
title_bar(s, "Embedding consolidation", "One (N, dim) dataset + has_embedding flag vs. per-UUID datasets")
add_figure(s, "opt_read_calls", Inches(0.6), Inches(1.3), Inches(7.2), Inches(4.3))
bullets(s, [
    ("Was: one dataset per embedding → one read per row", 0, LIGHT_GREY),
    ("Now: a single slice read for all rows in range", 0, LIGHT_GREY),
    ("Array full-scan: 10,005 → 8 HDF5 calls", 0, GREEN),
    ("Context read latency 3.8 → 2.2 ms", 0, GREEN),
    ("Cost: ≈10% slower array writes (shared-chunk RMW)", 0, RED),
    ("Not exercised by the live hook (no embeddings) ⁴", 0, YELLOW),
], Inches(8.0), Inches(1.4), Inches(4.9), Inches(5.4), head="Takeaways", size=14)

# ── 10. ORC ──────────────────────────────────────────────────────────────────
s = blank_slide(prs); fill_bg(s)
title_bar(s, "ORC as a log backend?", "Columnar archival format vs. live append")
add_figure(s, "opt_orc_cliff", Inches(0.6), Inches(1.4), Inches(6.6), Inches(4.1))
bullets(s, [
    ("ORC is write-once: a stripe is sealed on write", 0, LIGHT_GREY),
    ("Batch mode: huge t/s — but buffers ALL turns in RAM, nothing durable until close", 0, DIM_WHITE),
    ("Live mode = rewrite whole file per turn → O(N²)", 0, RED),
    ("Array-heavy: 132 → 55 t/s as N grows 50 → 500", 0, RED),
    ("Verdict: great for archival/export, poor for live logging", 0, YELLOW),
    ("The HDF5 schema keeps a tunable durability/throughput knob", 0, GREEN),
], Inches(7.4), Inches(1.4), Inches(5.5), Inches(5.4), head="Takeaways", size=14)

# ── Live hijack: what & how ──────────────────────────────────────────────────
s = blank_slide(prs); fill_bg(s)
title_bar(s, "Live session capture (the 'hijack')", "Incremental JSONL → HDF5 while you use Claude Code")
bullets(s, [
    ("Registered as Claude Code hooks on UserPromptSubmit and Stop", 0, LIGHT_GREY),
    ("Each event: read JSONL lines written since the last call (cursor stored in the HDF5 file)", 1, DIM_WHITE),
    ("Convert new user/assistant turns + tool calls to HDF5 rows", 1, DIM_WHITE),
    ("Match tool_use → tool_result across events via a pending-map attr", 1, DIM_WHITE),
    ("Never blocks Claude Code: all errors are swallowed, exit 0", 1, DIM_WHITE),
    ("Result: a queryable .h5 mirror of the session, updated live", 0, GREEN),
    ("Inspect with: `agentic-conversations-hdf5 inspect <file>` / `tail <file> <sid>`", 0, LIGHT_GREY),
], Inches(0.6), Inches(1.3), Inches(7.4), Inches(5.8), head="How it works")
add_rect(s, Inches(8.2), Inches(1.3), Inches(4.7), Inches(5.4), CODEBG)
add_text(s,
    "~/.claude/settings.json\n\n"
    '"hooks": {\n'
    '  "UserPromptSubmit": [{"hooks": [{\n'
    '    "type": "command",\n'
    '    "command": "python3 -m agentic_\n'
    '      conversations_hdf5.hooks.\n'
    '      live_session"}]}],\n'
    '  "Stop": [{"hooks": [{ ... }]}]\n'
    "}\n\n"
    "# one-time setup:\n"
    "$ agentic-conversations-hdf5 \\\n"
    "    setup-hook\n",
    Inches(8.35), Inches(1.45), Inches(4.4), Inches(5.1), size=11, color=LIGHT_GREY)

# ── 13. Enabling it / applicability ──────────────────────────────────────────
s = blank_slide(prs); fill_bg(s)
title_bar(s, "Turning it on — and what to expect", "From pip install to a live HDF5 session mirror")
bullets(s, [
    ("pip install agentic-conversations-hdf5", 0, GREEN),
    ("agentic-conversations-hdf5 setup-hook   (writes the two hook entries)", 0, LIGHT_GREY),
    ("Use Claude Code normally; .h5 files appear under ~/.claude/hdf5-sessions/", 0, LIGHT_GREY),
    ("agentic-conversations-hdf5 teardown-hook   to disable", 0, DIM_WHITE),
    ("Which optimizations matter here?", 0, YELLOW),
    ("libver, pre-allocation, zlib-ng, chunk defaults — all active", 1, GREEN),
    ("Batched buffer — little effect: few new lines per event ³", 1, DIM_WHITE),
    ("Embedding consolidation — not exercised: CC logs have no embeddings ⁴", 1, DIM_WHITE),
], Inches(0.6), Inches(1.3), Inches(12.2), Inches(5.8), head="Enable / disable", size=15)

# ── 14. Summary ──────────────────────────────────────────────────────────────
s = blank_slide(prs); fill_bg(s)
title_bar(s, "Summary")
bullets(s, [
    ("Write path is dramatically faster for bulk conversion (≈87× cumulative; batched buffer is the step change)", 0, LIGHT_GREY),
    ("Read path is leaner: embedding consolidation cuts a full scan from 10,005 HDF5 calls to 8", 0, LIGHT_GREY),
    ("Methodology is explicit: exact call counts, honest variance/profiler caveats", 0, LIGHT_GREY),
    ("ORC is archival-only; the HDF5 schema remains the right live-logging backend", 0, LIGHT_GREY),
    ("Live session capture works today via a Claude Code hook — pip-installable, one-command setup", 0, GREEN),
    ("Footnotes matter: the headline write speedups target bulk conversion; the live hook's wins are the always-on libver/pre-alloc/zlib-ng layer", 0, YELLOW),
], Inches(0.6), Inches(1.3), Inches(12.2), Inches(5.8), head="Where things stand", size=15)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"wrote {OUT.relative_to(ROOT)}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
